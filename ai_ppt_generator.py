import sys
import os
import configparser
import requests
import json
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, 
                             QLabel, QLineEdit, QPushButton, QTextEdit, QProgressBar, 
                             QFileDialog, QSlider, QMenuBar, QAction, QMessageBox)
from PyQt5.QtCore import Qt, QThread, pyqtSignal
from PyQt5.QtGui import QIcon, QFont

def get_config_path():
    if getattr(sys, 'frozen', False):
        # 如果是打包后的可执行文件
        application_path = os.path.dirname(sys.executable)
    else:
        # 如果是脚本运行
        application_path = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(application_path, 'config.ini')

# 读取配置文件
config = configparser.ConfigParser()
config_path = get_config_path()

if os.path.exists(config_path):
    config.read(config_path)
else:
    # 如果外部配置文件不存在，使用默认配置
    config.read(os.path.join(os.path.dirname(__file__), 'config.ini'))

OPENROUTER_API_KEY = config.get('API', 'OPENROUTER_API_KEY')
YOUR_SITE_URL = config.get('API', 'YOUR_SITE_URL')
YOUR_APP_NAME = config.get('API', 'YOUR_APP_NAME')

class PPTGeneratorThread(QThread):
    update_progress = pyqtSignal(int)
    update_status = pyqtSignal(str)
    finished = pyqtSignal(str, str)

    def __init__(self, topic, slide_count):
        super().__init__()
        self.topic = topic
        self.slide_count = slide_count

    def run(self):
        try:
            self.update_status.emit("正在生成PPT内容...")
            self.update_progress.emit(25)
            content = self.generate_ppt_content(self.topic, self.slide_count)
            
            self.update_status.emit("正在创建PPT文件...")
            self.update_progress.emit(75)
            output_file = f"{self.topic.replace(' ', '_')}_presentation.pptx"
            prs = self.create_ppt(content, output_file)
            
            self.update_progress.emit(100)
            self.finished.emit(content, output_file)
        except Exception as e:
            self.update_status.emit(f"错误: {str(e)}")

    def generate_ppt_content(self, topic, slide_count):
        response = requests.post(
            url="https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "HTTP-Referer": YOUR_SITE_URL,
                "X-Title": YOUR_APP_NAME,
            },
            data=json.dumps({
                "model": "openai/gpt-3.5-turbo",
                "messages": [
                    {"role": "system", "content": f"你是一个AI助手，负责创建PowerPoint演示文稿。请严格按照要求提供恰好{slide_count}张幻灯片的演示文稿内容，遵循以下结构和格式："},
                    {"role": "user", "content": f"""请创建一个关于"{topic}"的演示文稿，严格包含{slide_count}张幻灯片，遵循以下结构和格式：

1. 封面 (1张)
格式：
幻灯片1：封面
标题：[主标题]
副标题：[副标题]

2. 目录 (1张)
格式：
幻灯片2：目录
- [第一章节]
- [第二章节]
- ...

3. 内容页 ({slide_count - 4}张)
格式：
幻灯片[数字]：[章节标题]
- [要点1]
- [要点2]
- ...
[建议插入的图片描述]

4. 结论 (1张)
格式：
幻灯片[数字]：结论
- [总结要点1]
- [总结要点2]
- ...

5. 参考资料 (1张)
格式：
幻灯片[数字]：参考资料
- [参考资料1]
- [参考资料2]
- ...

请确保每张幻灯片的格式严格遵循上述要求。特别注意：
1. 每张幻灯片必须以"幻灯片[数字]："开头，其中[数字]必须是一个有效的整数。
2. 封面幻灯片的标题和副标题必须分别以"标题："和"副标题："开头。
3. 其他幻灯片的内容使用短横线列表。
4. 对于内容页，请在要点之后添加一行建议插入的图片描述。
5. 所有内容都应该使用中文。
6. 确保总幻灯片数量严格等于{slide_count}张，且幻灯片编号从1开始连续递增。
7. 根据幻灯片数量，合理分配内容，确保每张幻灯片的内容简洁明了，易于理解。"""}
                ]
            })
        )
        
        if response.status_code == 200:
            content = response.json()['choices'][0]['message']['content']
            return content
        else:
            raise Exception(f"错误: {response.status_code}, {response.text}")

    def create_ppt(self, content, output_file):
        prs = Presentation()
        
        slides = content.split('\n\n')
        
        for slide_content in slides:
            lines = slide_content.split('\n')
            if not lines:
                continue
            
            try:
                slide_info = lines[0].split('：')
                if len(slide_info) < 2:
                    continue
                
                slide_number_str = slide_info[0].split('幻灯片')
                if len(slide_number_str) < 2:
                    continue
                
                slide_number = int(slide_number_str[1])
                slide_title = slide_info[1]
            except (ValueError, IndexError) as e:
                print(f"Error processing slide: {e}")
                continue

            if slide_number == 1:  # 封面
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                # 查找标题和副标题
                title_text = ""
                subtitle_text = ""
                for line in lines[1:]:
                    if line.startswith("标题："):
                        title_text = line.split("：", 1)[1].strip()
                    elif line.startswith("副标题："):
                        subtitle_text = line.split("：", 1)[1].strip()
                
                title.text = title_text
                subtitle.text = subtitle_text
            elif slide_number == 2:  # 目录
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content = slide.placeholders[1]
                title.text = slide_title
                tf = content.text_frame
                
                # 计算内容的行数
                content_lines = lines[1:]
                line_count = len(content_lines)
                
                # 根据行数动态调整字体大小
                if line_count > 10:
                    font_size = max(8, int(18 - (line_count - 10) * 0.5))  # 最小字号为8
                else:
                    font_size = 18
                
                for line in content_lines:
                    p = tf.add_paragraph()
                    p.text = line.strip('- ')
                    p.level = 0
                    p.font.size = Pt(font_size)
            else:  # 内容页、结论和参考资料
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content = slide.placeholders[1]
                title.text = slide_title
                tf = content.text_frame
                for line in lines[1:]:
                    if line.startswith('[建议插入的图片描述]'):
                        p = tf.add_paragraph()
                        p.text = line.strip('[]')
                        p.font.size = Pt(10)
                        p.font.color.rgb = RGBColor(128, 128, 128)
                    else:
                        p = tf.add_paragraph()
                        p.text = line.strip('- ')
                        p.level = 0 if not line.startswith('  ') else 1

        return prs

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("AI PPT生成器")
        self.setGeometry(100, 100, 500, 600)
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f0f0f0;
            }
            QLabel {
                font-size: 14px;
                color: #333;
            }
            QLineEdit, QTextEdit {
                border: 1px solid #ccc;
                border-radius: 4px;
                padding: 5px;
                font-size: 14px;
            }
            QPushButton {
                background-color: #4CAF50;
                color: white;
                border: none;
                padding: 8px 16px;
                font-size: 14px;
                border-radius: 4px;
            }
            QPushButton:hover {
                background-color: #45a049;
            }
            QPushButton:disabled {
                background-color: #cccccc;
            }
            QSlider::groove:horizontal {
                border: 1px solid #999999;
                height: 8px;
                background: qlineargradient(x1:0, y1:0, x2:0, y2:1, stop:0 #B1B1B1, stop:1 #c4c4c4);
                margin: 2px 0;
            }
            QSlider::handle:horizontal {
                background: qlineargradient(x1:0, y1:0, x2:1, y2:1, stop:0 #b4b4b4, stop:1 #8f8f8f);
                border: 1px solid #5c5c5c;
                width: 18px;
                margin: -2px 0;
                border-radius: 3px;
            }
        """)
        
        # 创建菜单栏
        self.create_menu()

        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        layout = QVBoxLayout(central_widget)
        layout.setSpacing(20)
        layout.setContentsMargins(20, 20, 20, 20)

        # 主题输入
        topic_layout = QHBoxLayout()
        topic_label = QLabel("PPT主题:")
        self.topic_input = QLineEdit()
        self.topic_input.setPlaceholderText("输入PPT主题")
        topic_layout.addWidget(topic_label)
        topic_layout.addWidget(self.topic_input)
        layout.addLayout(topic_layout)

        # 幻灯片数量
        slide_count_layout = QVBoxLayout()
        slide_count_label = QLabel("幻灯片数量:")
        self.slide_count_slider = QSlider(Qt.Horizontal)
        self.slide_count_slider.setRange(5, 50)  # 修改这里，最大值改为50
        self.slide_count_slider.setValue(10)
        self.slide_count_slider.setTickPosition(QSlider.TicksBelow)
        self.slide_count_slider.setTickInterval(5)  # 修改刻度间隔为5
        self.slide_count_value = QLabel("10")
        self.slide_count_value.setAlignment(Qt.AlignCenter)
        self.slide_count_slider.valueChanged.connect(self.update_slide_count)
        slide_count_layout.addWidget(slide_count_label)
        slide_count_layout.addWidget(self.slide_count_slider)
        slide_count_layout.addWidget(self.slide_count_value)
        layout.addLayout(slide_count_layout)

        # 生成按钮
        self.generate_button = QPushButton("生成PPT")
        self.generate_button.clicked.connect(self.generate_ppt)
        layout.addWidget(self.generate_button)

        # 状态和进度
        self.status_label = QLabel("就绪")
        layout.addWidget(self.status_label)
        self.progress_bar = QProgressBar()
        self.progress_bar.setTextVisible(False)
        layout.addWidget(self.progress_bar)

        # 内容预览
        preview_label = QLabel("内容预览:")
        layout.addWidget(preview_label)
        self.content_preview = QTextEdit()
        self.content_preview.setReadOnly(True)
        self.content_preview.setMinimumHeight(200)
        layout.addWidget(self.content_preview)

        # 保存按钮
        self.save_button = QPushButton("导出PPT")
        self.save_button.clicked.connect(self.save_ppt)
        self.save_button.setEnabled(False)
        layout.addWidget(self.save_button)

        self.ppt_content = ""
        self.ppt_object = None
        self.generator_thread = None

    def create_menu(self):
        menubar = self.menuBar()

        # 文件菜单
        file_menu = menubar.addMenu('文件')
        new_action = QAction('新建', self)
        new_action.triggered.connect(self.new_file)
        file_menu.addAction(new_action)
        save_action = QAction('保存', self)
        save_action.triggered.connect(self.save_ppt)
        file_menu.addAction(save_action)
        exit_action = QAction('退出', self)
        exit_action.triggered.connect(self.close)
        file_menu.addAction(exit_action)

        # 编辑菜单
        edit_menu = menubar.addMenu('编辑')
        clear_action = QAction('清除内容', self)
        clear_action.triggered.connect(self.clear_content)
        edit_menu.addAction(clear_action)

        # 帮助菜单
        help_menu = menubar.addMenu('帮助')
        about_action = QAction('关于', self)
        about_action.triggered.connect(self.show_about)
        help_menu.addAction(about_action)

    def new_file(self):
        self.topic_input.clear()
        self.slide_count_slider.setValue(10)
        self.content_preview.clear()
        self.ppt_content = ""
        self.ppt_object = None
        self.save_button.setEnabled(False)
        self.status_label.setText("就绪")

    def clear_content(self):
        self.content_preview.clear()
        self.ppt_content = ""
        self.ppt_object = None
        self.save_button.setEnabled(False)
        self.status_label.setText("内容已清除")

    def show_about(self):
        QMessageBox.about(self, "关于", "AI PPT生成器\n\n版本 1.0\n\n使用AI技术自动生成PPT内容的工具。")

    def update_slide_count(self, value):
        self.slide_count_value.setText(str(value))

    def generate_ppt(self):
        topic = self.topic_input.text()
        if not topic:
            self.status_label.setText("请输入PPT主题")
            return

        slide_count = self.slide_count_slider.value()

        self.progress_bar.setValue(0)
        self.content_preview.clear()
        self.save_button.setEnabled(False)
        self.generate_button.setEnabled(False)

        if self.generator_thread is not None:
            self.generator_thread.quit()
            self.generator_thread.wait()

        self.generator_thread = PPTGeneratorThread(topic, slide_count)
        self.generator_thread.update_progress.connect(self.progress_bar.setValue)
        self.generator_thread.update_status.connect(self.status_label.setText)
        self.generator_thread.finished.connect(self.on_generation_finished)
        self.generator_thread.start()

    def on_generation_finished(self, content, output_file):
        self.ppt_content = content
        self.content_preview.setPlainText(content)
        self.status_label.setText("PPT内容生成完成")
        self.save_button.setEnabled(True)
        self.generate_button.setEnabled(True)
        self.generator_thread = None

    def save_ppt(self):
        if not self.ppt_content:
            return
        
        file_path, _ = QFileDialog.getSaveFileName(self, "保存PPT", "", "PowerPoint 文件 (*.pptx)")
        if file_path:
            try:
                ppt_generator = PPTGeneratorThread(self.topic_input.text(), self.slide_count_slider.value())
                prs = ppt_generator.create_ppt(self.ppt_content, file_path)
                prs.save(file_path)
                self.status_label.setText(f"PPT已保存: {file_path}")
            except Exception as e:
                self.status_label.setText(f"保存PPT时出错: {str(e)}")

    def closeEvent(self, event):
        if self.generator_thread is not None:
            self.generator_thread.quit()
            self.generator_thread.wait()
        event.accept()

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())